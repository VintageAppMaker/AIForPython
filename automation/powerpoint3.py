from pptx import Presentation
import requests
from io import BytesIO
from pptx.util import Inches, Pt

# 1. ppt 파일 만들기
prs = Presentation()

# 2. 이미지 가져와서 삽입하기
image_url = "https://chatsam.ai/_next/image?url=%2FChatSam%402x_003399.png&w=256&q=75"
response = requests.get(image_url)
image = BytesIO(response.content)

# 레이아웃이 없는 비어있는 슬라이드 추가하기
slide_layout = prs.slide_layouts[6]  # 레이아웃 없는 슬라이드 레이아웃 선택
slide = prs.slides.add_slide(slide_layout)

left = top = Inches(1)
slide.shapes.add_picture(image, left, top)

# 이미지 하단에 텍스트 추가하기
left = Inches(1)
top = Inches(4)
width = height = Inches(5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.add_paragraph()
p.text = "회사로고"
p.font.size = Pt(20)

# 3. ppt 파일 저장하기
ppt_file = "1.pptx"
prs.save(ppt_file)
