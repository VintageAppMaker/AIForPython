# 1차 에러발생 
# "Inches"를 찾지 못함.

# from pptx import Presentation
# import requests
# from io import BytesIO

# # 1. ppt 파일 만들기
# prs = Presentation()

# # 2. 이미지 가져와서 삽입하기
# image_url = "https://chatsam.ai/_next/image?url=%2FChatSam%402x_003399.png&w=256&q=75"
# response = requests.get(image_url)
# image = BytesIO(response.content)

# slide_layout = prs.slide_layouts[1]  # 두 번째 슬라이드 레이아웃 선택
# slide = prs.slides.add_slide(slide_layout)
# left = top = Inches(1)
# slide.shapes.add_picture(image, left, top)

# # 3. ppt 파일 저장하기
# ppt_file = "1.pptx"
# prs.save(ppt_file)

from pptx import Presentation
import requests
from io import BytesIO
from pptx.util import Inches

# 1. ppt 파일 만들기
prs = Presentation()

# 2. 이미지 가져와서 삽입하기
image_url = "https://chatsam.ai/_next/image?url=%2FChatSam%402x_003399.png&w=256&q=75"
response = requests.get(image_url)
image = BytesIO(response.content)

slide_layout = prs.slide_layouts[1]  # 두 번째 슬라이드 레이아웃 선택
slide = prs.slides.add_slide(slide_layout)
left = top = Inches(1)
slide.shapes.add_picture(image, left, top)

# 3. ppt 파일 저장하기
ppt_file = "1.pptx"
prs.save(ppt_file)
